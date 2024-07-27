import json
import logging

from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Pt

from vedanta.backend.apps.promptops import process_user_question
from vedanta.backend.knowledge_database.helper import get_knowledge_base, set_json_slide


async def generate_presentation(selected_theme="Theme1"):
    """
    Asynchronously generates a PowerPoint presentation based on the knowledge base.

    This function first retrieves the knowledge base using the `get_knowledge_base` function. If the knowledge base
    is empty, it raises an HTTPException with a status code of 400 and a detail message of "Please upload a PDF file
    first."

    The function then creates a prompt by combining the requirement for a PPTX presentation and the knowledge base.
    The prompt is a string that instructs the user to write a PowerPoint presentation in JSON format. The JSON should
    contain a 'slides' field, and each slide should have a 'title' and 'content' field. The content string of each
    slide should not exceed 45 words. If there is more to the topic, a new slide should be added.

    The function then calls the `process_user_question` function with the prompt as an argument to get a response.
    The response is converted to a string and printed.

    The function extracts the JSON string from the response by finding the index of the first occurrence of the
    substring '"s' and constructs the JSON string. The JSON string is then loaded into a Python dictionary using the
    `json.loads` function.

    The function calls the `set_json_slide` function with the JSON data to format the JSON data.

    The function then calls the `process_presentation` function with the JSON data to generate the PowerPoint
    presentation. The generated output file is returned.

    If any exception occurs during the execution of the function, an error message is logged and an HTTPException is
    raised with a status code of 500 and the error message as the detail.

    Returns:
        The generated PowerPoint presentation as a file.

    Raises:
        HTTPException: If the knowledge base is empty or if an error occurs during the execution of the function.
    """
    context = get_knowledge_base()
    if not context:
        raise HTTPException(status_code=400, detail="Please upload a PDF file first.")

    # Create the prompt by combining the requirement for PPTX and PDF context
    prompt = (
        f"Write a powerpoint presentation about this document in JSON. The JSON should contain 'slides'"
        f" and each slide should contain 'title' and 'content'. Make sure to not just use plain text"
        f"use bullet points and other things to make the text more engaging for the student."
        f"If there is more to the topic add a new slide. Do not be afraid to"
        f" divide things into 2-3 points for a slide."
    )

    try:
        print("prompt ab jayega")
        response_str = str(process_user_question(prompt))
        print("response kuch aaisa aaya")
        print(response_str)

        json_string = response_str.strip("<> \n```json")
        print(json_string)  # Extract the JSON string from the response
        # print("ln32")
        json_pptx = json.loads(json_string)
        print("formatting kuch aaise hui")
        print(json_pptx)
        set_json_slide(json_data=json_pptx)
        print("abb processing hogi")
        output_file = process_presentation(json_pptx, selected_theme)
        return output_file
    except Exception as e:
        logging.error(f"An error occurred during generate_presentation(): {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


theme_dict = {
    "Theme1": (0, 1),
    "Theme2": (11, 14),
    "Theme3": (0, 10),
    "Theme4": (0, 10),
}


def apply_bold_format(word, r, bolded_mode):
    """Applies bold formatting to words enclosed in double asterisks.

    This function checks if a word is enclosed in double asterisks ('**') and
    applies bold formatting to it within a PowerPoint text run.

    Args:
        word (str): The word to process for bold formatting.
        r (pptx.text.run.Run): The PowerPoint text run object to apply formatting to.
        bolded_mode (bool): A flag indicating whether the previous word was bolded.

    Returns:
        bool: True if the current word is bolded, False otherwise.
    """
    # Check if the word starts and ends with double asterisks
    if word.startswith("**") and word.endswith("**"):
        r.text = word[2:-2]  # Remove outer asterisks
        r.font.bold = True
    elif word.startswith("**"):
        r.text = word[2:] + " "  # Remove leading asterisks
        r.font.bold = True
        bolded_mode = True  # Switch to bolded mode
    elif word.endswith("**"):
        r.text = word[:-2] + " "  # Remove trailing asterisks
        r.font.bold = True
        bolded_mode = False  # Exit bolded mode
    elif bolded_mode:
        r.text = word + " "  # Apply bolded mode to other words
        r.font.bold = True
    else:
        r.text = word + " "  # Normal text, no bold
    return bolded_mode


def create_slide(prs, title, content, selected_theme):
    """Creates a new PowerPoint slide with the given title and content.

    Args:
        prs (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str): The title of the slide.
        content (str): The content of the slide.
        :param title:
        :param prs:
        :param content:
        :param selected_theme:
    """
    # Retrieve the theme selected by the user
    layout, text_placeholder = theme_dict.get(selected_theme, (11, 14))  # Default to (11, 14)

    # Create a new slide with appropriate layout
    slide_layout = prs.slide_layouts[layout]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the title text box
    title_shape = slide.shapes.title
    title_shape.text = title  # Assign the title

    # Use the content placeholder to set the content text
    content_shape = slide.placeholders[text_placeholder]
    content_text_frame = content_shape.text_frame
    content_text_frame.word_wrap = True  # Ensure proper word wrapping

    # Split content into lines based on newline character
    lines = content.replace("\\n", "\n").split("\n")  # Convert escaped newlines

    for line in lines:
        if line.startswith("##"):
            # Heading line
            p = content_text_frame.add_paragraph()
            p.text = line[2:].strip()  # Remove "##"
            p.font.size = Pt(20)  # Larger font for headings
            p.font.bold = True
        elif line.startswith("* "):
            # Bullet point
            p = content_text_frame.add_paragraph()
            line = line[1:].strip()
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)  # Apply bolding
            for run in p.runs:
                run.font.size = Pt(12)  # Smaller font for bullet points
        else:
            # Normal text
            p = content_text_frame.add_paragraph()
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)  # Apply bold if needed
                r.font.size = Pt(14)  # Standard font size for regular text


def process_presentation(json_data, selected_theme="Theme1"):
    """
    Process the given JSON data to generate a PowerPoint presentation.

    Args:
        json_data (dict or list): The JSON data representing the slides of the presentation.
            If a dictionary is provided, it should have a "slides" field containing a list of slide data.
            If a list is provided, it directly represents the slide data.
            Each slide data should be a dictionary with "title" and "content" fields.

    Returns:
        str: The path to the generated PowerPoint presentation file.

    Raises:
        HTTPException: If there is an error in the JSON data structure or during the presentation generation.
        :param json_data:
        :param selected_theme:
    """
    try:
        validate_json_data_structure(json_data)
        print("json valid tha")
        presentation = Presentation("./theme_pptx/" + selected_theme + ".pptx")
        slides = json_data.get("slides", []) if isinstance(json_data, dict) else json_data

        for slide_data in slides:
            if isinstance(slide_data, dict):
                title = slide_data.get("title", "")
                content = slide_data.get("content", "")
                content = content.replace("-", "")
                create_slide(presentation, title, content, selected_theme)

        output_file = "file.pptx"
        presentation.save(output_file)                                     # Save the presentation
        print("file toh save hogayi")

        return output_file
    except Exception as e:
        logging.error(f"An error occurred during process_presentation(): {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


def validate_json_data_structure(json_data):
    if not isinstance(json_data, (dict, list)):
        raise ValueError("Invalid JSON data format")

    if isinstance(json_data, dict):
        slides = json_data.get("slides")
        if slides is None:
            raise ValueError("Missing 'slides' field in slide data")
    else:
        slides = json_data

    for slide_data in slides:
        if not isinstance(slide_data, dict):
            raise ValueError("Invalid JSON data format")
        if "title" not in slide_data:
            raise ValueError("Missing 'title' field in slide data")
        if not isinstance(slide_data["title"], str):
            raise ValueError("'title' field must be a string")
        if "content" not in slide_data:
            raise ValueError("Missing 'content' field in slide data")
        if not isinstance(slide_data["content"], str):
            raise ValueError("'content' field must be a string")
